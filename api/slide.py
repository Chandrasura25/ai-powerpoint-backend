from fastapi import HTTPException, APIRouter
from fastapi.responses import FileResponse
from pydantic import BaseModel
import openai
from pptx import Presentation
import re
import os
from Config.config import Config
from pptx.util import Pt

router = APIRouter()

# OpenAI API Key
openai.api_key = Config.OPENAI_API_KEY

class PresentationRequest(BaseModel):
    topic: str
    num_slides: int = 5
    layout: str = "Varied"  # Default to Varied layout

@router.post("/generate-presentation")
async def generate_presentation(request: PresentationRequest):
    try:
        # Generate slide content
        slide_content = generate_slide_content(request.topic, request.num_slides)

        # Create PowerPoint presentation
        ppt = Presentation()

        for content in slide_content:
            # Choose layout based on request
            if request.layout == "Varied":
                slide = ppt.slides.add_slide(ppt.slide_layouts[8])  # Image & Text
            elif request.layout == "Text-Heavy":
                slide = ppt.slides.add_slide(ppt.slide_layouts[1])  # Title & Content
            else:  # Image-Focused
                slide = ppt.slides.add_slide(ppt.slide_layouts[5])  # Title & Image

            # Add title (ensure it appears at the top and is bigger)
            if slide.shapes.title:
                slide.shapes.title.text = content["title"]
                slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(32) 
            # Add bullet points only if there's a content placeholder
            if request.layout in ["Varied", "Text-Heavy"]:
                try:
                    text_placeholder = slide.placeholders[1].text_frame  # Content placeholder
                    text_placeholder.clear()
                    
                    for point in content["bullet_points"]:
                        p = text_placeholder.add_paragraph()
                        p.text = point
                        p.space_after = 10  # Adds spacing between bullets
                except IndexError:
                    pass  # Skip if content placeholder does not exist

            # Add image if applicable
            if request.layout in ["Varied", "Image-Focused"]:
                add_image_to_slide(slide)

        # Save the PowerPoint file
        file_path = f"{request.topic}.pptx"
        ppt.save(file_path)

        return FileResponse(file_path, filename=f"{request.topic}.pptx")

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


def generate_slide_content(topic: str, num_slides: int):
    """Fetch slide content using OpenAI API."""
    try:
        prompt = f"Generate an outline for a PowerPoint presentation on '{topic}' with {num_slides} slides. Each slide should have a title and 3-5 bullet points."
        response = openai.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "system", "content": "You are a helpful assistant."},
                      {"role": "user", "content": prompt}],
            max_tokens=500,
        )

        raw_text = response.choices[0].message.content.strip()

        slides = []
        slide_texts = raw_text.split("\n\n")  # Split by double newlines

        for slide_text in slide_texts:
            lines = slide_text.strip().split("\n")
            if not lines:
                continue

            title = re.sub(r"[*]+", "", lines[0]).strip()  # Remove `**`
            bullet_points = [re.sub(r"[*-]+", "", point).strip() for point in lines[1:] if point.strip()]

            slides.append({"title": title, "bullet_points": bullet_points})

        return slides

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


def add_image_to_slide(slide):
    """Add an image to the slide in the correct position."""
    try:
        # Placeholder image path (replace with dynamic images if needed)
        image_path = os.path.join(os.path.dirname(__file__), "placeholder.png")

        # Define position (adjust as needed to fit layout)
        left, top, width, height = 400, 150, 300, 250  # Adjusted for better placement
        slide.shapes.add_picture(image_path, left, top, width, height)

    except Exception:
        pass  # Ignore image errors if any occur
